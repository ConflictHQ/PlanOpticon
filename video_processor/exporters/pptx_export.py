"""Generate PPTX slide decks from knowledge graph data.

Uses python-pptx for slide generation. Falls back gracefully if not installed.
No LLM required — pure template-based generation from KG data.
"""

import logging
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)


def _add_title_slide(prs, title: str, subtitle: str):
    """Add a title slide."""
    from pptx.util import Pt

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    body.text = subtitle
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)


def _add_content_slide(prs, title: str, bullets: List[str]):
    """Add a slide with bullet points."""
    from pptx.util import Pt

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
            for run in tf.paragraphs[0].runs:
                run.font.size = Pt(14)
        else:
            p = tf.add_paragraph()
            p.text = bullet
            for run in p.runs:
                run.font.size = Pt(14)


def _add_table_slide(prs, title: str, headers: List[str], rows: List[List[str]]):
    """Add a slide with a table."""
    from pptx.util import Emu, Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide.shapes.title.text = title

    num_rows = len(rows) + 1
    num_cols = len(headers)

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    row_height = Emu(int(Inches(0.35)))
    height = row_height * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)


def _add_image_slide(prs, title: str, image_path: Path):
    """Add a slide with an embedded image."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide.shapes.title.text = title

    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    slide.shapes.add_picture(str(image_path), left, top, width=width)


def generate_pptx(
    kg_data: dict,
    output_path: Path,
    title: Optional[str] = None,
    diagrams_dir: Optional[Path] = None,
) -> Path:
    """Generate a PPTX slide deck from knowledge graph data.

    Args:
        kg_data: Knowledge graph dict with 'nodes' and 'relationships'.
        output_path: Path to write the PPTX file.
        title: Optional presentation title.
        diagrams_dir: Optional directory containing diagram images to embed.

    Returns:
        Path to the generated PPTX.

    Raises:
        ImportError: If python-pptx is not installed.
    """
    from pptx import Presentation

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    nodes = kg_data.get("nodes", [])
    rels = kg_data.get("relationships", [])

    report_title = title or "Knowledge Graph"
    now = datetime.now().strftime("%Y-%m-%d %H:%M")

    # Title slide
    _add_title_slide(
        prs,
        report_title,
        f"Generated {now}\n{len(nodes)} entities \u2022 {len(rels)} relationships",
    )

    # Overview slide
    by_type: Dict[str, list] = {}
    for n in nodes:
        t = n.get("type", "concept")
        by_type.setdefault(t, []).append(n)

    overview_bullets = [f"{len(nodes)} entities across {len(by_type)} types"]
    for etype, elist in sorted(by_type.items(), key=lambda x: -len(x[1])):
        examples = ", ".join(e.get("name", "") for e in elist[:3])
        overview_bullets.append(f"{etype.title()} ({len(elist)}): {examples}")
    _add_content_slide(prs, "Overview", overview_bullets)

    # Key entities slide
    degree: Dict[str, int] = {}
    for r in rels:
        degree[r.get("source", "")] = degree.get(r.get("source", ""), 0) + 1
        degree[r.get("target", "")] = degree.get(r.get("target", ""), 0) + 1

    top = sorted(degree.items(), key=lambda x: -x[1])[:10]
    if top:
        _add_table_slide(
            prs,
            "Key Entities",
            ["Entity", "Connections"],
            [[name, str(deg)] for name, deg in top],
        )

    # Diagram slides
    if diagrams_dir and diagrams_dir.exists():
        pngs = sorted(diagrams_dir.glob("*.png"))
        for i, png in enumerate(pngs):
            _add_image_slide(prs, f"Diagram {i + 1}", png)

    # Relationship types slide
    rel_types: Dict[str, int] = {}
    for r in rels:
        rt = r.get("type", "related_to")
        rel_types[rt] = rel_types.get(rt, 0) + 1

    if rel_types:
        _add_table_slide(
            prs,
            "Relationship Types",
            ["Type", "Count"],
            [[rt, str(c)] for rt, c in sorted(rel_types.items(), key=lambda x: -x[1])],
        )

    # Entity detail slides (batched, max 12 per slide)
    batch_size = 12
    for batch_start in range(0, len(nodes), batch_size):
        batch = nodes[batch_start : batch_start + batch_size]
        bullets = []
        for node in batch:
            name = node.get("name", "")
            etype = node.get("type", "concept")
            descs = node.get("descriptions", [])
            desc = descs[0][:80] if descs else ""
            bullets.append(f"{name} ({etype}): {desc}")

        slide_num = batch_start // batch_size + 1
        total_pages = (len(nodes) + batch_size - 1) // batch_size
        page_label = f" ({slide_num}/{total_pages})" if total_pages > 1 else ""
        _add_content_slide(prs, f"Entities{page_label}", bullets)

    prs.save(str(output_path))
    logger.info(f"Generated PPTX: {output_path}")
    return output_path
