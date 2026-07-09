"""Tests for PDF and PPTX exporters."""

import pytest

from video_processor.exporters.pdf_export import generate_pdf
from video_processor.exporters.pptx_export import generate_pptx


def _sample_kg():
    """Return a sample knowledge graph dict for testing."""
    return {
        "nodes": [
            {"name": "Python", "type": "technology", "descriptions": ["A programming language"]},
            {"name": "Django", "type": "technology", "descriptions": ["A web framework"]},
            {"name": "Alice", "type": "person", "descriptions": ["Software engineer"]},
            {"name": "Bob", "type": "person", "descriptions": ["Product manager"]},
            {"name": "Acme Corp", "type": "organization", "descriptions": ["A tech company"]},
        ],
        "relationships": [
            {"source": "Alice", "target": "Python", "type": "uses"},
            {"source": "Alice", "target": "Bob", "type": "works_with"},
            {"source": "Django", "target": "Python", "type": "built_on"},
            {"source": "Alice", "target": "Acme Corp", "type": "employed_by"},
        ],
    }


def _empty_kg():
    return {"nodes": [], "relationships": []}


class TestPDFExport:
    @pytest.fixture(autouse=True)
    def _check_reportlab(self):
        pytest.importorskip("reportlab")

    def test_generate_pdf(self, tmp_path):
        out = tmp_path / "report.pdf"
        result = generate_pdf(_sample_kg(), out, title="Test Report")
        assert result == out
        assert out.exists()
        assert out.stat().st_size > 0

    def test_generate_pdf_empty_kg(self, tmp_path):
        out = tmp_path / "empty.pdf"
        result = generate_pdf(_empty_kg(), out)
        assert result == out
        assert out.exists()

    def test_generate_pdf_creates_parent_dirs(self, tmp_path):
        out = tmp_path / "sub" / "dir" / "report.pdf"
        result = generate_pdf(_sample_kg(), out)
        assert result == out
        assert out.exists()

    def test_generate_pdf_default_title(self, tmp_path):
        out = tmp_path / "default.pdf"
        generate_pdf(_sample_kg(), out)
        assert out.exists()

    def test_generate_pdf_with_diagrams_dir(self, tmp_path):
        diag_dir = tmp_path / "diagrams"
        diag_dir.mkdir()
        out = tmp_path / "report.pdf"
        # No PNGs in dir — should still work
        result = generate_pdf(_sample_kg(), out, diagrams_dir=diag_dir)
        assert result == out

    def test_generate_pdf_no_reportlab(self, tmp_path, monkeypatch):
        """Verify ImportError propagates when reportlab is missing."""
        import builtins

        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name.startswith("reportlab"):
                raise ImportError("No module named 'reportlab'")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)
        with pytest.raises(ImportError):
            generate_pdf(_sample_kg(), tmp_path / "fail.pdf")


class TestPPTXExport:
    @pytest.fixture(autouse=True)
    def _check_pptx(self):
        pytest.importorskip("pptx")

    def test_generate_pptx(self, tmp_path):
        out = tmp_path / "slides.pptx"
        result = generate_pptx(_sample_kg(), out, title="Test Deck")
        assert result == out
        assert out.exists()
        assert out.stat().st_size > 0

    def test_generate_pptx_empty_kg(self, tmp_path):
        out = tmp_path / "empty.pptx"
        result = generate_pptx(_empty_kg(), out)
        assert result == out
        assert out.exists()

    def test_generate_pptx_creates_parent_dirs(self, tmp_path):
        out = tmp_path / "sub" / "dir" / "slides.pptx"
        result = generate_pptx(_sample_kg(), out)
        assert result == out
        assert out.exists()

    def test_generate_pptx_with_diagrams_dir(self, tmp_path):
        diag_dir = tmp_path / "diagrams"
        diag_dir.mkdir()
        out = tmp_path / "slides.pptx"
        result = generate_pptx(_sample_kg(), out, diagrams_dir=diag_dir)
        assert result == out

    def test_pptx_slide_count(self, tmp_path):
        """Verify expected number of slides are created."""
        from pptx import Presentation

        out = tmp_path / "count.pptx"
        generate_pptx(_sample_kg(), out)
        prs = Presentation(str(out))
        # Title + Overview + Key Entities + Rel Types + 1 entity batch = 5
        assert len(prs.slides) == 5

    def test_pptx_many_entities_batched(self, tmp_path):
        """Entities are batched into multiple slides when >12."""
        from pptx import Presentation

        kg = {
            "nodes": [
                {"name": f"Entity{i}", "type": "concept", "descriptions": [f"desc {i}"]}
                for i in range(25)
            ],
            "relationships": [],
        }
        out = tmp_path / "many.pptx"
        generate_pptx(kg, out)
        prs = Presentation(str(out))
        # Title + Overview + 3 entity batches (12 + 12 + 1) = 5
        # No Key Entities or Rel Types slides (no relationships)
        assert len(prs.slides) == 5

    def test_generate_pptx_no_pptx(self, tmp_path, monkeypatch):
        """Verify ImportError propagates when python-pptx is missing."""
        import builtins

        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name.startswith("pptx"):
                raise ImportError("No module named 'pptx'")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)
        with pytest.raises(ImportError):
            generate_pptx(_sample_kg(), tmp_path / "fail.pptx")


class TestConflictKGExport:
    def test_canonical_shape(self):
        from video_processor.exporters.conflict_kg import FORMAT_ID, to_conflict_kg

        data = to_conflict_kg(_sample_kg())
        assert data["format"] == FORMAT_ID == "conflict-kg/v1"
        assert len(data["nodes"]) == 5
        assert len(data["edges"]) == 4

        ids = [n["id"] for n in data["nodes"]]
        assert len(ids) == len(set(ids))
        assert ids == [n["name"].lower() for n in data["nodes"]]

        python = next(n for n in data["nodes"] if n["id"] == "python")
        assert python["name"] == "Python"
        assert python["type"] == "technology"
        assert python["props"]["descriptions"] == ["A programming language"]

    def test_edges_reference_node_ids(self):
        from video_processor.exporters.conflict_kg import to_conflict_kg

        data = to_conflict_kg(_sample_kg())
        node_ids = {n["id"] for n in data["nodes"]}
        for edge in data["edges"]:
            assert edge["source"] in node_ids
            assert edge["target"] in node_ids

    def test_edge_props_dropped_when_absent(self):
        from video_processor.exporters.conflict_kg import to_conflict_kg

        kg = {
            "nodes": [{"name": "A", "type": "concept"}, {"name": "B", "type": "concept"}],
            "relationships": [
                {"source": "A", "target": "B", "type": "related_to", "timestamp": 5.0},
                {"source": "B", "target": "A", "type": "related_to", "content_source": None},
            ],
        }
        data = to_conflict_kg(kg)
        assert data["edges"][0]["props"] == {"timestamp": 5.0}
        assert data["edges"][1]["props"] == {}

    def test_json_writer(self, tmp_path):
        import json

        from video_processor.exporters.conflict_kg import to_conflict_kg, write_conflict_kg_json

        out = tmp_path / "conflict_kg.json"
        write_conflict_kg_json(_sample_kg(), out)
        assert json.loads(out.read_text()) == to_conflict_kg(_sample_kg())

    def test_sqlite_writer_schema_and_rows(self, tmp_path):
        import json
        import sqlite3

        from video_processor.exporters.conflict_kg import write_conflict_kg_sqlite

        out = tmp_path / "conflict_kg.db"
        write_conflict_kg_sqlite(_sample_kg(), out)

        conn = sqlite3.connect(out)
        try:
            tables = {
                r[0]
                for r in conn.execute(
                    "SELECT name FROM sqlite_master WHERE type='table'"
                ).fetchall()
            }
            assert tables == {"nodes", "edges"}
            indexes = {
                r[0]
                for r in conn.execute(
                    "SELECT name FROM sqlite_master WHERE type='index' AND name LIKE 'idx%'"
                ).fetchall()
            }
            assert indexes == {"idx_edges_source", "idx_edges_target"}

            nodes = conn.execute("SELECT id, name, type, props FROM nodes ORDER BY id").fetchall()
            assert len(nodes) == 5
            assert all(json.loads(props) is not None for *_rest, props in nodes)

            edges = conn.execute("SELECT source, target, type FROM edges").fetchall()
            assert len(edges) == 4
            node_ids = {n[0] for n in nodes}
            assert all(src in node_ids and tgt in node_ids for src, tgt, _ in edges)
        finally:
            conn.close()

    def test_sqlite_matches_json_encoding(self, tmp_path):
        import json
        import sqlite3

        from video_processor.exporters.conflict_kg import (
            to_conflict_kg,
            write_conflict_kg_sqlite,
        )

        out = tmp_path / "conflict_kg.db"
        write_conflict_kg_sqlite(_sample_kg(), out)
        canonical = to_conflict_kg(_sample_kg())

        conn = sqlite3.connect(out)
        try:
            nodes = [
                {"id": i, "name": n, "type": t, "props": json.loads(p)}
                for i, n, t, p in conn.execute("SELECT id, name, type, props FROM nodes")
            ]
            edges = [
                {"source": s, "target": t, "type": ty, "props": json.loads(p)}
                for s, t, ty, p in conn.execute("SELECT source, target, type, props FROM edges")
            ]
        finally:
            conn.close()
        assert nodes == canonical["nodes"]
        assert edges == canonical["edges"]

    def test_cli_export_conflict_kg(self, tmp_path):
        import json

        from click.testing import CliRunner

        from video_processor.cli.commands import cli
        from video_processor.integrators.graph_store import SQLiteStore

        db = tmp_path / "knowledge_graph.db"
        store = SQLiteStore(db)
        store.merge_entity("Alice", "person", ["Engineer"])
        store.merge_entity("Python", "technology", [])
        store.add_relationship("Alice", "Python", "uses")

        runner = CliRunner()
        out_json = tmp_path / "out.json"
        result = runner.invoke(cli, ["export", "conflict-kg", str(db), "-o", str(out_json)])
        assert result.exit_code == 0, result.output
        data = json.loads(out_json.read_text())
        assert data["format"] == "conflict-kg/v1"
        assert {n["id"] for n in data["nodes"]} == {"alice", "python"}
        assert data["edges"][0] == {
            "source": "alice",
            "target": "python",
            "type": "uses",
            "props": {},
        }

        out_db = tmp_path / "out.db"
        result = runner.invoke(
            cli, ["export", "conflict-kg", str(db), "--sqlite", "-o", str(out_db)]
        )
        assert result.exit_code == 0, result.output
        assert out_db.exists()
